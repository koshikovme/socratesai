# -*- coding: utf-8 -*-
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
TEMPLATE = next(BASE.glob("*.pptx"))
ASSETS = BASE / "generated_presentation_assets"
OUTPUT = BASE / "SocratesAI_Research_Work_Presentation_Koshikov_AA.pptx"

FIG_PIPELINE = ASSETS / "report_fig_01.png"
FIG_SUBMISSION = ASSETS / "report_fig_02.png"
FIG_ABLATION = ASSETS / "report_fig_03.png"
FIG_UI = ASSETS / "report_fig_04.png"
FIG_CONFUSION = ASSETS / "report_fig_05.png"


NAVY = RGBColor(12, 35, 64)
TEAL = RGBColor(0, 153, 166)
CYAN = RGBColor(42, 190, 210)
GREEN = RGBColor(34, 153, 84)
AMBER = RGBColor(236, 157, 45)
RED = RGBColor(192, 57, 43)
PURPLE = RGBColor(104, 89, 180)
INK = RGBColor(33, 43, 54)
MUTED = RGBColor(94, 108, 132)
LIGHT_BG = RGBColor(246, 248, 250)
LIGHT_BLUE = RGBColor(232, 248, 250)
LIGHT_GREEN = RGBColor(236, 248, 240)
LIGHT_AMBER = RGBColor(254, 246, 232)
LIGHT_PURPLE = RGBColor(242, 240, 252)
WHITE = RGBColor(255, 255, 255)
MID_GRAY = RGBColor(214, 221, 230)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.lstrip("#")
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


def clear_existing_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def set_run(run, size=18, color=INK, bold=False, italic=False):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_text(shape, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    text_frame.margin_top = Inches(0.04)
    text_frame.margin_bottom = Inches(0.04)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)


def add_textbox(slide, x, y, w, h, text="", size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box, text, size=size, color=color, bold=bold, align=align)
    return box


def add_title(slide, title, idx=None):
    title_shape = slide.shapes.title
    if title_shape is not None:
        text_frame = title_shape.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        set_run(run, size=24, color=NAVY, bold=True)
    else:
        add_textbox(slide, 0.65, 0.25, 11.8, 0.55, title, size=24, color=NAVY, bold=True)
    if idx is not None:
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(12.0),
            Inches(0.28),
            Inches(0.7),
            Inches(0.38),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = TEAL
        pill.line.color.rgb = TEAL
        set_text(pill, f"{idx:02d}", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.02), Inches(11.65), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.color.rgb = TEAL


def add_footer(slide, number):
    add_textbox(
        slide,
        0.72,
        7.02,
        8.5,
        0.22,
        "Koshikov A.A. | SocratesAI | Master's student research work #2",
        size=7.8,
        color=MUTED,
    )
    add_textbox(slide, 12.0, 7.02, 0.55, 0.22, str(number), size=7.8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, items, size=16, color=INK, bullet=True, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        if bullet:
            p.text = f"• {item}"
            for run in p.runs:
                set_run(run, size=size, color=color)
        else:
            run = p.add_run()
            run.text = item
            set_run(run, size=size, color=color)
    return box


def add_card(slide, x, y, w, h, title, body="", fill=WHITE, accent=TEAL, title_size=15, body_size=11.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = MID_GRAY
    shape.line.width = Pt(0.75)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    text = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.10), Inches(w - 0.28), Inches(h - 0.16))
    tf = text.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=title_size, color=NAVY, bold=True)
    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.line_spacing = 1.0
        r2 = p2.add_run()
        r2.text = body
        set_run(r2, size=body_size, color=INK)
    return shape


def add_metric(slide, x, y, w, h, value, label, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_GRAY
    text = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.08), Inches(w - 0.16), Inches(h - 0.14))
    tf = text.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    set_run(r, size=21, color=accent, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = label
    set_run(r2, size=8.8, color=MUTED)
    return shape


def add_image_fit(slide, image_path, x, y, w, h, border=True, bg=True):
    if bg:
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = WHITE
        bg_shape.line.color.rgb = MID_GRAY
        bg_shape.line.width = Pt(0.75)
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    aspect = img_w / img_h
    box_aspect = w / h
    pad = 0.12
    iw = w - pad * 2
    ih = h - pad * 2
    if aspect >= box_aspect:
        final_w = iw
        final_h = iw / aspect
    else:
        final_h = ih
        final_w = ih * aspect
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    pic = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))
    if border:
        pic.line.color.rgb = MID_GRAY
        pic.line.width = Pt(0.4)
    return pic


def add_flow_step(slide, x, y, w, h, text, fill, accent=TEAL, size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.0)
    set_text(shape, text, size=size, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, x, y, w=0.35, h=0.18, color=TEAL):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.color.rgb = color
    return arrow


def add_timeline(slide, milestones):
    y = 3.25
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y), Inches(11.05), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_GRAY
    line.line.color.rgb = MID_GRAY
    xs = [1.0, 2.85, 4.75, 6.65, 8.65, 10.75]
    for i, (label, date, status, color) in enumerate(milestones):
        x = xs[i]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y - 0.14), Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        add_textbox(slide, x - 0.28, y - 0.78, 1.25, 0.32, date, size=8.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - 0.62, y + 0.28, 1.75, 0.75, label, size=8.7, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - 0.55, y + 0.96, 1.6, 0.42, status, size=7.6, color=MUTED, align=PP_ALIGN.CENTER)


def make_deck():
    prs = Presentation(str(TEMPLATE))
    clear_existing_slides(prs)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    set_text(
        title,
        "Research and Implementation of Virtual Mentor Models in Programming",
        size=28,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    subtitle = slide.placeholders[1]
    set_text(
        subtitle,
        "Real-time personalized feedback for first-year students\n\n"
        "Master's student research work #2, 2025-2026 academic year\n"
        "Alimzhan Koshikov | 7M06105 Computer Science and Engineering\n"
        "Supervisor: Tutkyshbayeva Shyryn, PhD | Astana IT University",
        size=15,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        1.35,
        6.28,
        10.7,
        0.46,
        "Dissertation topic: Исследование и реализация моделей виртуального наставника по программированию "
        "с персонализированной обратной связью студентам первого курса в реальном времени",
        size=8.8,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "1. Relevance of the Research / Актуальность исследования", 2)
    add_card(
        slide,
        0.82,
        1.48,
        3.65,
        1.18,
        "Delayed help in CS1",
        "Beginners often need feedback at the exact moment they are stuck, not after grading.",
        fill=LIGHT_BLUE,
        accent=TEAL,
    )
    add_card(
        slide,
        0.82,
        2.92,
        3.65,
        1.18,
        "Static tools are limited",
        "Compilers and tests report failure, but rarely explain the student's misconception.",
        fill=LIGHT_AMBER,
        accent=AMBER,
    )
    add_card(
        slide,
        0.82,
        4.36,
        3.65,
        1.18,
        "Unrestricted LLMs are risky",
        "Open-ended assistants can over-explain or reveal near-solutions without pedagogical control.",
        fill=LIGHT_PURPLE,
        accent=PURPLE,
    )
    add_textbox(slide, 5.2, 1.46, 6.65, 0.38, "Research gap", size=19, color=NAVY, bold=True)
    add_bullets(
        slide,
        5.22,
        2.02,
        6.5,
        1.8,
        [
            "Need feedback that is immediate, personalized, and bounded by learning goals.",
            "The mentor must decide how to intervene before generating the wording.",
            "The key technical layer is a feedback-state / policy decision model.",
        ],
        size=15,
    )
    add_flow_step(slide, 5.25, 4.25, 1.8, 0.75, "Code state", LIGHT_BG, NAVY)
    add_arrow(slide, 7.15, 4.53)
    add_flow_step(slide, 7.65, 4.25, 2.2, 0.75, "Feedback state", LIGHT_BLUE, TEAL)
    add_arrow(slide, 10.02, 4.53)
    add_flow_step(slide, 10.52, 4.25, 1.65, 0.75, "Bounded hint", LIGHT_GREEN, GREEN)
    add_textbox(
        slide,
        5.22,
        5.45,
        6.7,
        0.72,
        "Core idea: feedback generation should be controlled by an explicit pedagogical policy, "
        "not by an unrestricted language model.",
        size=15,
        color=NAVY,
        bold=True,
    )
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2. Purpose and Objectives / Цель и задачи исследования", 3)
    add_card(
        slide,
        0.82,
        1.38,
        11.35,
        0.96,
        "Purpose",
        "Develop and justify the architecture of a programming virtual mentor that provides real-time personalized feedback for first-year students.",
        fill=LIGHT_BLUE,
        accent=TEAL,
        title_size=16,
        body_size=13.2,
    )
    objectives = [
        ("01", "Analyze existing CS1 feedback systems and identify limitations of static assessment."),
        ("02", "Study LLM and code-analysis methods for diagnosing beginner programming errors."),
        ("03", "Define an adaptive feedback model: state prediction, policy action, bounded hint text."),
        ("04", "Implement a SocratesAI MVP with analyzer, policy selector, feedback generator, and logs."),
        ("05", "Evaluate the ML decision layer and prepare classroom validation for 2026-2027."),
    ]
    x_positions = [0.9, 3.15, 5.4, 7.65, 9.9]
    for x, (num, txt) in zip(x_positions, objectives):
        add_metric(slide, x, 2.85, 1.72, 1.05, num, "objective", accent=TEAL)
        add_textbox(slide, x - 0.1, 4.03, 1.96, 1.2, txt, size=9.8, color=INK, align=PP_ALIGN.CENTER)
    add_card(
        slide,
        1.35,
        5.62,
        10.25,
        0.72,
        "Research question",
        "Can a virtual mentor select an appropriate pedagogical feedback state before hint generation, using code, run results, and student-context signals?",
        fill=WHITE,
        accent=NAVY,
        title_size=14,
        body_size=12.3,
    )
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "3. Object, Subject and Scientific Novelty / Объект, предмет, новизна", 4)
    add_card(
        slide,
        0.82,
        1.42,
        5.25,
        1.12,
        "Object of research",
        "Learning programming in CS1/CS0 courses with automated and AI-assisted feedback.",
        fill=LIGHT_BG,
        accent=NAVY,
        title_size=15,
        body_size=12.4,
    )
    add_card(
        slide,
        6.45,
        1.42,
        5.25,
        1.12,
        "Subject of research",
        "Models and architecture for real-time personalized feedback in a programming virtual mentor.",
        fill=LIGHT_BG,
        accent=TEAL,
        title_size=15,
        body_size=12.4,
    )
    add_textbox(slide, 0.88, 2.98, 5.2, 0.3, "Scientific novelty", size=17, color=NAVY, bold=True)
    add_bullets(
        slide,
        0.9,
        3.38,
        5.35,
        2.0,
        [
            "Explicit decision layer between code analysis and LLM/text feedback.",
            "Feedback state is treated as a supervised prediction problem.",
            "Interaction logs store action, confidence, latency, and outcome labels for later evaluation.",
        ],
        size=13.2,
    )
    add_textbox(slide, 6.48, 2.98, 5.2, 0.3, "Contributions in this period", size=17, color=NAVY, bold=True)
    add_bullets(
        slide,
        6.5,
        3.38,
        5.35,
        2.0,
        [
            "Research direction refined after SIST 2026 reviews.",
            "New IEEE manuscript submitted with ML feedback-state methodology.",
            "SocratesAI MVP implemented and connected to the ML policy path.",
        ],
        size=13.2,
    )
    add_card(
        slide,
        1.05,
        5.82,
        10.95,
        0.56,
        "Boundary",
        "At this stage the evidence supports technical feasibility and action-selection validity; classroom learning-gain testing is planned next.",
        fill=LIGHT_AMBER,
        accent=AMBER,
        title_size=12.5,
        body_size=10.5,
    )
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "4. Methods and Initial Data / Методы и исходные данные", 5)
    methods = [
        ("Literature review", "CS1 learning difficulties, automated feedback, AI tutors, LLM guardrails"),
        ("Supervised ML", "Feedback-state prediction with problem-level holdout and ablation study"),
        ("Software prototyping", "Spring Boot backend, Vue frontend, PostgreSQL, FastAPI ML policy service"),
        ("Evaluation design", "Offline Codeforces evaluation plus SocratesAI interaction logs for future CS1 testing"),
    ]
    for i, (t, b) in enumerate(methods):
        add_card(
            slide,
            0.8 + (i % 2) * 5.75,
            1.36 + (i // 2) * 1.58,
            5.25,
            1.08,
            t,
            b,
            fill=[LIGHT_BLUE, LIGHT_GREEN, LIGHT_PURPLE, LIGHT_AMBER][i],
            accent=[TEAL, GREEN, PURPLE, AMBER][i],
            title_size=14.2,
            body_size=11.3,
        )
    add_textbox(slide, 0.85, 4.68, 3.2, 0.32, "Main data used", size=17, color=NAVY, bold=True)
    add_metric(slide, 0.88, 5.16, 2.05, 1.05, "178,414", "Python submissions", accent=TEAL)
    add_metric(slide, 3.18, 5.16, 2.05, 1.05, "872", "beginner Codeforces problems", accent=GREEN)
    add_metric(slide, 5.48, 5.16, 2.05, 1.05, "5", "feedback states", accent=PURPLE)
    add_metric(slide, 7.78, 5.16, 2.05, 1.05, "12", "MVP benchmark tasks", accent=AMBER)
    add_metric(slide, 10.08, 5.16, 2.05, 1.05, "2026-2027", "classroom validation", accent=RED)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "5. ML Decision Layer / Слой принятия решения", 6)
    add_image_fit(slide, FIG_PIPELINE, 0.88, 1.35, 4.55, 5.15)
    add_textbox(slide, 5.75, 1.42, 6.15, 0.38, "From submission to mentor action", size=18, color=NAVY, bold=True)
    add_bullets(
        slide,
        5.78,
        1.98,
        6.05,
        2.25,
        [
            "Beginner submissions are grouped into feedback categories from judge/run outcomes.",
            "Train/test split is performed by problem to reduce leakage across similar attempts.",
            "The predictor estimates the feedback state before the mentor writes the hint.",
            "SocratesAI uses the state and confidence to guide bounded feedback generation.",
        ],
        size=13.6,
    )
    add_flow_step(slide, 5.88, 4.82, 1.5, 0.68, "Analyzer", LIGHT_BG, NAVY, size=11)
    add_arrow(slide, 7.55, 5.08, 0.32, 0.15)
    add_flow_step(slide, 8.02, 4.82, 1.65, 0.68, "Policy", LIGHT_BLUE, TEAL, size=11)
    add_arrow(slide, 9.84, 5.08, 0.32, 0.15)
    add_flow_step(slide, 10.28, 4.82, 1.7, 0.68, "Feedback", LIGHT_GREEN, GREEN, size=11)
    add_card(
        slide,
        5.72,
        5.86,
        6.22,
        0.55,
        "Design principle",
        "The LLM or template generates wording only after the policy action is selected.",
        fill=LIGHT_AMBER,
        accent=AMBER,
        title_size=11.8,
        body_size=9.8,
    )
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "6. SocratesAI MVP Architecture / Архитектура прототипа", 7)
    y1, y2, y3 = 1.58, 3.22, 4.9
    add_flow_step(slide, 0.78, y2, 1.55, 0.72, "Student UI\nVue 3", LIGHT_BLUE, TEAL, 10.6)
    add_arrow(slide, 2.48, y2 + 0.27, 0.38, 0.18)
    add_flow_step(slide, 3.0, y2, 1.95, 0.72, "REST / WebSocket\nSpring Boot", LIGHT_BG, NAVY, 10.4)
    add_arrow(slide, 5.12, y2 + 0.27, 0.38, 0.18)
    add_flow_step(slide, 5.65, y1, 1.8, 0.72, "Analyzer\nsignals", LIGHT_GREEN, GREEN, 10.4)
    add_flow_step(slide, 5.65, y2, 1.8, 0.72, "Session\ncontext", LIGHT_PURPLE, PURPLE, 10.4)
    add_flow_step(slide, 5.65, y3, 1.8, 0.72, "Task + run\nresults", LIGHT_AMBER, AMBER, 10.4)
    add_arrow(slide, 7.62, y2 + 0.27, 0.38, 0.18)
    add_flow_step(slide, 8.15, y2, 1.82, 0.72, "Policy selector\nRule / ML", LIGHT_BLUE, TEAL, 10.4)
    add_arrow(slide, 10.15, y2 + 0.27, 0.38, 0.18)
    add_flow_step(slide, 10.68, y1, 1.8, 0.72, "Feedback\nTemplate / LLM", LIGHT_GREEN, GREEN, 10.2)
    add_flow_step(slide, 10.68, y3, 1.8, 0.72, "interaction_logs\nPostgreSQL", LIGHT_BG, NAVY, 10.2)
    add_textbox(slide, 0.95, 5.98, 11.2, 0.4, "Implementation evidence from source code", size=16, color=NAVY, bold=True)
    add_bullets(
        slide,
        1.0,
        6.36,
        11.0,
        0.5,
        [
            "MentorWorkflowService: analyzer -> context -> policy -> feedback; MlPolicySelector calls FastAPI; InteractionLog persists action, confidence, latency, and outcome fields."
        ],
        size=10.5,
        bullet=False,
    )
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "7. Feedback Policy and Actions / Политика обратной связи", 8)
    add_textbox(slide, 0.85, 1.32, 11.3, 0.42, "The model does not only say whether code is correct; it selects the pedagogical type of help.", size=16, color=NAVY, bold=True)
    actions = [
        ("CODE_HIGHLIGHT", "Point to a suspicious code region when location matters more than explanation.", LIGHT_BLUE, TEAL),
        ("CONCEPTUAL_HINT", "Give a short conceptual nudge without exposing the full solution.", LIGHT_GREEN, GREEN),
        ("GUIDING_QUESTION", "Ask a Socratic question that makes the student inspect their own logic.", LIGHT_PURPLE, PURPLE),
        ("NO_FEEDBACK", "Avoid interruption when the evidence is weak or the code state is locally acceptable.", LIGHT_AMBER, AMBER),
    ]
    for i, (t, b, fill, accent) in enumerate(actions):
        add_card(slide, 0.9 + (i % 2) * 5.65, 2.05 + (i // 2) * 1.55, 5.15, 1.05, t, b, fill=fill, accent=accent, title_size=14.2, body_size=11.2)
    add_card(
        slide,
        1.05,
        5.65,
        10.95,
        0.68,
        "Research-state taxonomy in the manuscript",
        "accepted, semantic debugging, execution safety, efficiency review, syntax repair. The runtime policy maps evidence to an action that controls feedback wording.",
        fill=WHITE,
        accent=NAVY,
        title_size=12.5,
        body_size=10.4,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "8. Empirical Results / Результаты ML-эксперимента", 9)
    add_metric(slide, 0.82, 1.35, 1.95, 0.98, "0.085", "majority baseline macro F1", accent=RED)
    add_metric(slide, 2.97, 1.35, 1.95, 0.98, "0.426", "source-code model macro F1", accent=TEAL)
    add_metric(slide, 5.12, 1.35, 1.95, 0.98, "0.814", "code + run model macro F1", accent=GREEN)
    add_metric(slide, 7.27, 1.35, 1.95, 0.98, "872", "held-out problem setting", accent=PURPLE)
    add_metric(slide, 9.42, 1.35, 2.38, 0.98, "178,414", "submissions evaluated", accent=AMBER)
    add_image_fit(slide, FIG_ABLATION, 0.98, 2.65, 6.45, 3.72)
    add_textbox(slide, 7.82, 2.72, 3.95, 0.32, "Interpretation", size=17, color=NAVY, bold=True)
    add_bullets(
        slide,
        7.82,
        3.18,
        4.15,
        2.15,
        [
            "Source code contains useful early signals for personalized mentoring.",
            "Execution results substantially improve reliable state prediction.",
            "Problem-level holdout makes the evaluation stricter than random splitting.",
        ],
        size=12.6,
    )
    add_card(
        slide,
        7.8,
        5.62,
        4.1,
        0.72,
        "Boundary",
        "These results validate the decision layer, not final learning gain.",
        fill=LIGHT_AMBER,
        accent=AMBER,
        title_size=11.5,
        body_size=9.7,
    )
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "9. Error Analysis / Анализ ошибок модели", 10)
    add_image_fit(slide, FIG_CONFUSION, 0.9, 1.36, 5.55, 5.18)
    add_textbox(slide, 6.85, 1.42, 5.1, 0.36, "What the confusion matrix shows", size=18, color=NAVY, bold=True)
    add_bullets(
        slide,
        6.87,
        1.98,
        4.95,
        2.05,
        [
            "Syntax repair is easier to separate than semantic/execution/efficiency states.",
            "Some accepted and semantic cases overlap when only source code is available.",
            "The model needs richer run-result and history features for stable decisions.",
        ],
        size=13.0,
    )
    add_card(
        slide,
        6.82,
        4.52,
        5.0,
        1.15,
        "Implication for SocratesAI",
        "The MVP stores feedback action, confidence, suspicious region, latency, and outcome labels so future classroom data can improve the decision layer.",
        fill=LIGHT_BLUE,
        accent=TEAL,
        title_size=13.4,
        body_size=10.8,
    )
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "10. SocratesAI MVP Implementation / Реализация MVP", 11)
    add_image_fit(slide, FIG_UI, 0.82, 1.32, 6.45, 3.38)
    add_textbox(slide, 7.62, 1.38, 4.35, 0.36, "Implemented modules", size=18, color=NAVY, bold=True)
    add_bullets(
        slide,
        7.62,
        1.92,
        4.25,
        1.82,
        [
            "Spring Boot Java 21 backend, Vue frontend, PostgreSQL storage.",
            "Analyzer, session context, rule/ML policy, and bounded feedback generation.",
            "REST and WebSocket paths for real-time mentor feedback.",
        ],
        size=12.6,
    )
    add_textbox(slide, 0.95, 5.02, 3.55, 0.34, "Project-side validation", size=17, color=NAVY, bold=True)
    add_metric(slide, 0.95, 5.48, 1.78, 0.82, "360", "controlled events", accent=TEAL)
    add_metric(slide, 2.95, 5.48, 1.78, 0.82, "0", "runtime errors", accent=GREEN)
    add_metric(slide, 4.95, 5.48, 1.78, 0.82, "95.83%", "agreement", accent=PURPLE)
    add_metric(slide, 6.95, 5.48, 1.78, 0.82, "95.80%", "macro F1", accent=AMBER)
    add_card(
        slide,
        8.95,
        5.3,
        2.95,
        1.02,
        "Interpretation",
        "Technical feasibility and action-selection validity; classroom pilot still required.",
        fill=LIGHT_AMBER,
        accent=AMBER,
        title_size=11.6,
        body_size=9.2,
    )
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "11. IMWP Implementation and Publication / Результаты ИПРМ", 12)
    milestones = [
        ("Literature review and theoretical basis", "Jan-Mar 2026", "completed", GREEN),
        ("SocratesAI concept and architecture", "Mar-Apr 2026", "completed", GREEN),
        ("SIST 2026 manuscript reviewed", "Mar 2026", "not accepted; feedback used", AMBER),
        ("Technical redirection to ML decision layer", "Apr-Jun 2026", "completed", GREEN),
        ("IEEE student conference manuscript", "Jun 2026", "submitted", TEAL),
        ("MVP and initial experiments", "May-Jun 2026", "completed; pilot planned", GREEN),
    ]
    add_timeline(slide, milestones)
    add_card(
        slide,
        0.95,
        1.42,
        5.15,
        0.78,
        "Publication output",
        '"A Machine Learning Model for Real-Time Personalized Feedback in Introductory Programming" - submitted to the 2026 IEEE 3rd International Student Conference on Digital Generation.',
        fill=LIGHT_BLUE,
        accent=TEAL,
        title_size=13,
        body_size=9.9,
    )
    add_card(
        slide,
        6.45,
        1.42,
        5.15,
        0.78,
        "Research internship",
        "Astana IT University, Kazakhstan; 09.03.2026-30.05.2026, recorded in the Individual Master's Work Plan.",
        fill=LIGHT_GREEN,
        accent=GREEN,
        title_size=13,
        body_size=9.9,
    )
    add_bullets(
        slide,
        1.05,
        5.55,
        10.9,
        0.8,
        [
            "The reporting period moved the dissertation from literature review and architecture planning to an implemented MVP plus empirical ML evaluation."
        ],
        size=12.8,
        color=NAVY,
        bullet=False,
    )
    add_footer(slide, 12)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Conclusion / Заключение", 13)
    add_card(
        slide,
        0.9,
        1.45,
        3.42,
        3.1,
        "What has been achieved",
        "The research direction was corrected toward a measurable technical layer. The MVP now connects code analysis, policy selection, bounded feedback, and interaction logging.",
        fill=LIGHT_BLUE,
        accent=TEAL,
        title_size=15,
        body_size=12.2,
    )
    add_card(
        slide,
        4.95,
        1.45,
        3.42,
        3.1,
        "Main scientific result",
        "Feedback-state prediction is feasible on a problem-level holdout: source code already improves strongly over baseline, and code-plus-run signals reach the best result.",
        fill=LIGHT_GREEN,
        accent=GREEN,
        title_size=15,
        body_size=12.2,
    )
    add_card(
        slide,
        9.0,
        1.45,
        3.0,
        3.1,
        "Next stage",
        "Collect target-domain SocratesAI logs, run CS1 classroom or historical-data validation, refine the taxonomy, and complete the implementation/results chapters.",
        fill=LIGHT_AMBER,
        accent=AMBER,
        title_size=15,
        body_size=12.2,
    )
    add_textbox(
        slide,
        1.45,
        5.45,
        10.4,
        0.62,
        "Dissertation defense deadline: 2027. The current work provides the technical foundation; the planned classroom study will test learning outcomes.",
        size=16,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 13)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = make_deck()
    print(output)
